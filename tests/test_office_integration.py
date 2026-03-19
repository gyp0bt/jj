"""Officeプラグイン統合テスト: PPTX/XLSXパーサー・エクスポーター

python-pptx / openpyxl がインストールされている場合のみ実行。

[READMEへ戻る](../README.md)
"""

from __future__ import annotations

import tempfile
from pathlib import Path
from unittest.mock import MagicMock

import pytest

# ================================================================
# パーサーテスト
# ================================================================


class TestPptxMetadataParser:
    """PptxMetadataParser のテスト"""

    @pytest.fixture
    def pptx_lib(self):
        return pytest.importorskip("pptx")

    def test_parser_registration(self, pptx_lib):
        """PptxMetadataParser が自動登録される"""
        from services.parse.connectors.office.pptx_parser import PptxMetadataParser

        assert PptxMetadataParser.priority == 35

    def test_extract_metadata_from_fixture(self, pptx_lib):
        """テストフィクスチャの .pptx からメタデータが抽出できる"""
        import zipfile

        fixture_pptx = Path("tests/fixtures/graph_test1/reports/260205_構造解析_idx1.pptx")
        if not fixture_pptx.exists() or not zipfile.is_zipfile(fixture_pptx):
            pytest.skip("有効なPPTXフィクスチャが見つかりません")

        from pptx import Presentation

        from services.parse.connectors.office.pptx_parser import _extract_pptx_metadata

        metadata = _extract_pptx_metadata(fixture_pptx, Presentation)
        assert "pptx.slide_count" in metadata
        assert isinstance(metadata["pptx.slide_count"], int)
        assert metadata["pptx.slide_count"] >= 0

    def test_extract_metadata_with_created_pptx(self, pptx_lib):
        """新規作成した .pptx からメタデータが正しく抽出できる"""
        from pptx import Presentation
        from pptx.util import Inches

        from services.parse.connectors.office.pptx_parser import _extract_pptx_metadata

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # スライド2枚追加
        layout = prs.slide_layouts[0]
        slide1 = prs.slides.add_slide(layout)
        if slide1.shapes.title:
            slide1.shapes.title.text = "テストスライド1"
        slide2 = prs.slides.add_slide(layout)
        if slide2.shapes.title:
            slide2.shapes.title.text = "テストスライド2"

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            prs.save(f.name)
            tmp_path = Path(f.name)

        try:
            metadata = _extract_pptx_metadata(tmp_path, Presentation)
            assert metadata["pptx.slide_count"] == 2
            assert "pptx.slide_width_cm" in metadata
            assert "pptx.slide_height_cm" in metadata
        finally:
            tmp_path.unlink()

    def test_parser_apply_skips_non_pptx(self, pptx_lib):
        """PPTX以外のフォーマットはスキップされる"""
        from services.parse.connectors.office.pptx_parser import PptxMetadataParser

        node = MagicMock()
        node.format = "csv"
        node.properties = {"path": "test.csv"}

        graph = MagicMock()
        graph.nodes = [node]

        original_props = dict(node.properties)
        parser = PptxMetadataParser()
        result = parser.apply(graph)
        assert result == graph
        # propertiesが変更されていないこと
        assert node.properties == original_props


class TestXlsxMetadataParser:
    """XlsxMetadataParser のテスト"""

    @pytest.fixture
    def openpyxl_lib(self):
        return pytest.importorskip("openpyxl")

    def test_parser_registration(self, openpyxl_lib):
        """XlsxMetadataParser が自動登録される"""
        from services.parse.connectors.office.xlsx_parser import XlsxMetadataParser

        assert XlsxMetadataParser.priority == 35

    def test_extract_metadata_from_created_xlsx(self, openpyxl_lib):
        """新規作成した .xlsx からメタデータが正しく抽出できる"""
        from openpyxl import Workbook, load_workbook

        from services.parse.connectors.office.xlsx_parser import _extract_xlsx_metadata

        wb = Workbook()
        ws1 = wb.active
        ws1.title = "Data"
        ws1.append(["col1", "col2", "col3"])
        ws1.append([1, 2, 3])
        ws1.append([4, 5, 6])

        ws2 = wb.create_sheet("Summary")
        ws2.append(["name", "value"])
        ws2.append(["test", 100])

        with tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False) as f:
            wb.save(f.name)
            tmp_path = Path(f.name)

        try:
            metadata = _extract_xlsx_metadata(tmp_path, load_workbook)
            assert metadata["xlsx.sheet_count"] == 2
            assert metadata["xlsx.sheet_names"] == ["Data", "Summary"]
            assert "xlsx.total_rows" in metadata
            assert "xlsx.sheet_dimensions" in metadata
            assert metadata["xlsx.sheet_dimensions"]["Data"]["columns"] == 3
        finally:
            tmp_path.unlink()

    def test_parser_apply_skips_non_xlsx(self, openpyxl_lib):
        """XLSX以外のフォーマットはスキップされる"""
        from services.parse.connectors.office.xlsx_parser import XlsxMetadataParser

        node = MagicMock()
        node.format = "csv"
        node.properties = {"path": "test.csv"}

        graph = MagicMock()
        graph.nodes = [node]

        original_props = dict(node.properties)
        parser = XlsxMetadataParser()
        result = parser.apply(graph)
        assert result == graph
        assert node.properties == original_props


# ================================================================
# Excelエクスポーターテスト
# ================================================================


class TestExcelExporter:
    """Excel書式付きエクスポーターのテスト"""

    @pytest.fixture
    def openpyxl_lib(self):
        return pytest.importorskip("openpyxl")

    def test_exporter_registration(self, openpyxl_lib):
        """ExcelExporter が自動登録される"""
        from services.export.connectors.excel_export import ExcelExporter

        assert ExcelExporter.format == "xlsx"
        assert ExcelExporter.priority == 12

    def test_export_table_to_excel(self, openpyxl_lib):
        """DataFrameをExcelに書式付きでエクスポートできる"""
        pd = pytest.importorskip("pandas")
        from services.export.connectors.excel_export import export_table_to_excel

        df = pd.DataFrame({"名前": ["テスト1", "テスト2"], "値": [100, 200], "備考": ["OK", "NG"]})

        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = export_table_to_excel(df, Path(tmpdir), filename_prefix="test")
            assert output_path.exists()
            assert output_path.suffix == ".xlsx"

            # 書式確認
            from openpyxl import load_workbook

            wb = load_workbook(str(output_path))
            ws = wb.active
            assert ws.title == "Data"
            # ヘッダー行の書式確認
            header_cell = ws.cell(row=1, column=1)
            assert header_cell.value == "名前"
            assert header_cell.font.bold is True
            assert header_cell.font.color.rgb == "00FFFFFF"
            # データ行
            assert ws.cell(row=2, column=1).value == "テスト1"
            assert ws.cell(row=2, column=2).value == 100

    def test_export_table_to_excel_bytes(self, openpyxl_lib):
        """DataFrameをExcelバイト列に変換できる"""
        pd = pytest.importorskip("pandas")
        from services.export.connectors.excel_export import export_table_to_excel_bytes

        df = pd.DataFrame({"col1": [1, 2, 3], "col2": ["a", "b", "c"]})
        result = export_table_to_excel_bytes(df)
        assert isinstance(result, bytes)
        assert len(result) > 0
        # XLSXマジックバイト（PKヘッダー）
        assert result[:2] == b"PK"

    def test_export_array_data_to_excel(self, openpyxl_lib):
        """配列データをExcelにエクスポートできる"""
        from services.export.connectors.excel_export import export_array_data_to_excel

        array_data = [
            {
                "name": "Node_1",
                "x": [0.0, 1.0, 2.0, 3.0],
                "y": [0.0, 1.0, 4.0, 9.0],
                "x_label": "Time",
                "y_label": "Displacement",
                "props": {"material": "Steel", "section": "S1"},
            },
            {
                "name": "Node_2",
                "x": [0.0, 1.0, 2.0],
                "y": [10.0, 20.0, 30.0],
                "x_label": "Time",
                "y_label": "Force",
                "props": {"material": "Aluminum"},
            },
        ]

        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = export_array_data_to_excel(array_data, Path(tmpdir))
            assert output_path.exists()

            from openpyxl import load_workbook

            wb = load_workbook(str(output_path))
            # Summary シートが先頭
            assert wb.sheetnames[0] == "Summary"
            # データシートが2つ
            assert len(wb.sheetnames) == 3

            # Summary シートの内容確認
            ws_summary = wb["Summary"]
            assert ws_summary.cell(row=1, column=1).value == "No."
            assert ws_summary.cell(row=2, column=2).value == "Node_1"

    def test_sanitize_sheet_name(self, openpyxl_lib):
        """シート名のサニタイズが正しく動作する"""
        from services.export.connectors.excel_export import _sanitize_sheet_name

        assert _sanitize_sheet_name("Normal", 1) == "Normal"
        assert _sanitize_sheet_name("Test[1]", 1) == "Test1"
        assert _sanitize_sheet_name("A" * 50, 1) == "A" * 25 + "_1"
        assert _sanitize_sheet_name("", 1) == "Sheet_1"

    def test_auto_column_width(self, openpyxl_lib):
        """列幅自動調整が正しく動作する"""
        from openpyxl import Workbook

        from services.export.connectors.excel_export import _auto_column_width

        wb = Workbook()
        ws = wb.active
        ws.cell(row=1, column=1, value="Short")
        ws.cell(row=1, column=2, value="これは長い日本語テキストです")

        _auto_column_width(ws)

        # 日本語列は英語列より幅が広いはず
        col_a_width = ws.column_dimensions["A"].width
        col_b_width = ws.column_dimensions["B"].width
        assert col_b_width > col_a_width


# ================================================================
# PPTXエクスポーターテスト
# ================================================================


class TestPptxExporter:
    """PPTX エクスポーターのテスト"""

    @pytest.fixture
    def pptx_lib(self):
        return pytest.importorskip("pptx")

    def test_exporter_registration(self, pptx_lib):
        """PptxExporter が自動登録される"""
        from services.export.connectors.pptx_export import PptxExporter

        assert PptxExporter.format == "pptx"
        assert PptxExporter.priority == 13

    def test_export_gallery_to_pptx(self, pptx_lib):
        """ギャラリー画像からPPTXを生成できる"""
        pytest.importorskip("PIL")
        # テスト用画像を作成
        from PIL import Image

        from services.export.connectors.pptx_export import export_gallery_to_pptx

        with tempfile.TemporaryDirectory() as tmpdir:
            images = []
            for i in range(5):
                img = Image.new("RGB", (200, 150), color=(i * 50, 100, 200))
                img_path = Path(tmpdir) / f"test_{i}.png"
                img.save(str(img_path))
                images.append(img_path)

            prs = export_gallery_to_pptx(images, cols=3, rows=2, title="テスト")
            assert len(prs.slides) == 1  # 5画像 < 3x2=6 → 1スライド

    def test_export_gallery_multiple_slides(self, pptx_lib):
        """画像数がグリッドを超える場合に複数スライドに分割される"""
        pytest.importorskip("PIL")
        from PIL import Image

        from services.export.connectors.pptx_export import export_gallery_to_pptx

        with tempfile.TemporaryDirectory() as tmpdir:
            images = []
            for i in range(8):
                img = Image.new("RGB", (100, 100), color=(200, i * 30, 100))
                img_path = Path(tmpdir) / f"test_{i}.png"
                img.save(str(img_path))
                images.append(img_path)

            prs = export_gallery_to_pptx(images, cols=2, rows=2, title="多スライド")
            assert len(prs.slides) == 2  # 8画像 / (2x2) = 2スライド

    def test_export_gallery_to_pptx_bytes(self, pptx_lib):
        """PPTXバイト列が正しく生成される"""
        pytest.importorskip("PIL")
        from PIL import Image

        from services.export.connectors.pptx_export import export_gallery_to_pptx_bytes

        with tempfile.TemporaryDirectory() as tmpdir:
            img_path = Path(tmpdir) / "test.png"
            img = Image.new("RGB", (100, 100), color="red")
            img.save(str(img_path))

            result = export_gallery_to_pptx_bytes([img_path])
            assert isinstance(result, bytes)
            assert result[:2] == b"PK"

    def test_export_plots_to_pptx(self, pptx_lib):
        """プロット画像からPPTXを生成できる"""
        pytest.importorskip("PIL")
        from PIL import Image

        from services.export.connectors.pptx_export import export_plots_to_pptx

        with tempfile.TemporaryDirectory() as tmpdir:
            plot_images = []
            for i in range(3):
                img = Image.new("RGB", (1920, 1080), color=(100, 200, i * 80))
                img_path = Path(tmpdir) / f"plot_{i}.png"
                img.save(str(img_path))
                plot_images.append({"path": img_path, "title": f"Plot {i + 1}"})

            prs = export_plots_to_pptx(plot_images)
            assert len(prs.slides) == 3  # 1図/スライド

    def test_export_gallery_no_valid_images_raises(self, pptx_lib):
        """有効な画像がない場合はValueErrorが発生する"""
        from services.export.connectors.pptx_export import export_gallery_to_pptx

        with pytest.raises(ValueError, match="有効な画像ファイル"):
            export_gallery_to_pptx([Path("/nonexistent/image.png")])


# ================================================================
# プラグイン登録テスト
# ================================================================


class TestOfficePluginRegistration:
    """Officeプラグインの登録テスト"""

    def test_plugin_module_imports(self):
        """Officeプラグインモジュールがインポートできる"""
        import services.plugins.office  # noqa: F401

    def test_register_is_idempotent(self):
        """register()は複数回呼んでも安全"""
        from services.plugins.office import register

        register()
        register()  # 2回目も例外なし
