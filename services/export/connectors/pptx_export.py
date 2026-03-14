"""PPTXエクスポーター: グラフデータ・画像をPowerPointファイルに出力

仕様書 windows-integration.md Phase W-3/W-4 に基づく実装:
- W-3: ギャラリーグリッド → PPTスライド（python-pptx使用、クロスプラットフォーム）
- W-4: プロット画像 → PPTスライド（1図/スライド）

python-pptxを使用し、新規.pptxファイルを生成する。
Win32 COM（アクティブなPowerPointへの貼り付け）は別途Windows専用モジュールで対応。

[READMEへ戻る](../../../README.md)
"""

from __future__ import annotations

import io
from datetime import datetime
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Emu, Inches, Pt

from jj_types import GraphModel
from services.export import AbstractExporter

# ================================================================
# 定数（スライドレイアウト・マージン）
# ================================================================

# スライドサイズ（ワイドスクリーン: 13.333 x 7.5 inch）
SLIDE_WIDTH_INCH = 13.333
SLIDE_HEIGHT_INCH = 7.5

# マージン
MARGIN_INCH = 0.5
TITLE_HEIGHT_INCH = 0.8
CELL_PADDING_INCH = 0.1

# 使用可能領域
USABLE_WIDTH_INCH = SLIDE_WIDTH_INCH - 2 * MARGIN_INCH
USABLE_HEIGHT_INCH = SLIDE_HEIGHT_INCH - 2 * MARGIN_INCH - TITLE_HEIGHT_INCH

# スライドレイアウトインデックス
LAYOUT_BLANK = 6
LAYOUT_TITLE_ONLY = 5


# ================================================================
# エクスポーター
# ================================================================


class PptxExporter(AbstractExporter):
    """PowerPoint形式エクスポーター

    ギャラリー画像やプロットをPowerPointファイルとして出力する。
    """

    format = "pptx"
    priority = 13

    def export(self, graph: GraphModel, **kwargs: Any) -> dict[str, Any]:
        """PPTXエクスポート

        kwargs:
            project_root: Path - プロジェクトルート
            images: list[Path] - 画像ファイルパスリスト
            cols: int - グリッド列数（デフォルト: 3）
            rows: int - グリッド行数（デフォルト: 2）
            title: str - プレゼンテーションタイトル
            output_file: str | None - 出力ファイル名
        """
        project_root = kwargs.get("project_root", Path("."))
        images = kwargs.get("images", [])
        cols = kwargs.get("cols", 3)
        rows = kwargs.get("rows", 2)
        title = kwargs.get("title", "Gallery Export")
        output_file = kwargs.get("output_file")

        if not images:
            raise ValueError("画像ファイルが指定されていません。")

        prs = export_gallery_to_pptx(
            images=[Path(p) for p in images],
            cols=cols,
            rows=rows,
            title=title,
        )

        if output_file is None:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            output_file = f"gallery_{timestamp}.pptx"
        output_path = project_root / ".j2" / "exports" / output_file
        output_path.parent.mkdir(parents=True, exist_ok=True)

        prs.save(str(output_path))
        return {"output_path": output_path, "count": len(images)}

    def format_cli_result(self, result: dict[str, Any], project_root: Path) -> str:
        output_path = result.get("output_path")
        count = result.get("count", 0)
        if output_path:
            try:
                rel = Path(output_path).relative_to(project_root)
            except ValueError:
                rel = output_path
            return f"PPTXエクスポート完了: {rel} ({count}画像)"
        return f"PPTXエクスポート完了 ({count}画像)"


# ================================================================
# ギャラリー → PPTスライド（W-3）
# ================================================================


def export_gallery_to_pptx(
    images: list[Path],
    cols: int = 3,
    rows: int = 2,
    title: str = "Gallery",
    slide_title_prefix: str = "",
) -> Presentation:
    """ギャラリー画像をPPTスライドにグリッド配置した新規プレゼンテーションを生成

    スライドサイズ（13.333 x 7.5 inch）からマージンを引いた領域を
    cols x rows に分割し、各セルに画像を配置する。
    画像数が cols * rows を超える場合は複数スライドに分割。

    Args:
        images: 画像ファイルパスリスト
        cols: グリッド列数
        rows: グリッド行数
        title: プレゼンテーションタイトル
        slide_title_prefix: スライドタイトルの接頭辞

    Returns:
        python-pptx Presentation オブジェクト
    """
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH_INCH)
    prs.slide_height = Inches(SLIDE_HEIGHT_INCH)

    images_per_slide = cols * rows
    existing_images = [img for img in images if img.exists()]

    if not existing_images:
        raise ValueError("有効な画像ファイルが見つかりません。")

    # スライド分割
    for slide_idx in range(0, len(existing_images), images_per_slide):
        slide_images = existing_images[slide_idx : slide_idx + images_per_slide]
        slide_num = slide_idx // images_per_slide + 1

        slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_BLANK])

        # スライドタイトル
        slide_title = f"{slide_title_prefix}{title}"
        if len(existing_images) > images_per_slide:
            total_slides = (len(existing_images) + images_per_slide - 1) // images_per_slide
            slide_title = f"{slide_title} ({slide_num}/{total_slides})"

        _add_slide_title(slide, slide_title)

        # グリッドに画像を配置
        _place_images_in_grid(slide, slide_images, cols, rows)

    return prs


def export_gallery_to_pptx_bytes(
    images: list[Path],
    cols: int = 3,
    rows: int = 2,
    title: str = "Gallery",
) -> bytes:
    """ギャラリー画像をPPTXバイト列に変換する（Streamlitダウンロード用）

    Args:
        images: 画像ファイルパスリスト
        cols: グリッド列数
        rows: グリッド行数
        title: プレゼンテーションタイトル

    Returns:
        PPTXファイルのバイト列
    """
    prs = export_gallery_to_pptx(images, cols, rows, title)
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()


# ================================================================
# プロット → PPTスライド（W-4）
# ================================================================


def export_plots_to_pptx(
    plot_images: list[dict[str, Any]],
    title: str = "Plots",
) -> Presentation:
    """プロット画像をPPTスライドに配置した新規プレゼンテーションを生成

    1図/スライド、全幅レイアウト。

    Args:
        plot_images: [{"path": Path, "title": str}, ...]
        title: プレゼンテーションタイトル

    Returns:
        python-pptx Presentation オブジェクト
    """
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH_INCH)
    prs.slide_height = Inches(SLIDE_HEIGHT_INCH)

    for item in plot_images:
        img_path = Path(item["path"])
        if not img_path.exists():
            continue

        slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_BLANK])

        # スライドタイトル
        plot_title = item.get("title", img_path.stem)
        _add_slide_title(slide, plot_title)

        # 画像を全幅で配置
        left = Inches(MARGIN_INCH)
        top = Inches(MARGIN_INCH + TITLE_HEIGHT_INCH)
        max_width = Inches(USABLE_WIDTH_INCH)
        max_height = Inches(USABLE_HEIGHT_INCH)

        _add_image_fit(slide, img_path, left, top, max_width, max_height)

    return prs


def export_plots_to_pptx_bytes(
    plot_images: list[dict[str, Any]],
    title: str = "Plots",
) -> bytes:
    """プロット画像をPPTXバイト列に変換する（Streamlitダウンロード用）

    Args:
        plot_images: [{"path": Path, "title": str}, ...]
        title: プレゼンテーションタイトル

    Returns:
        PPTXファイルのバイト列
    """
    prs = export_plots_to_pptx(plot_images, title)
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()


# ================================================================
# Plotly/Matplotlib → PNG → PPTX ユーティリティ
# ================================================================


def plotly_fig_to_pptx_bytes(
    figs: list[Any],
    titles: list[str] | None = None,
) -> bytes:
    """Plotly figureリストをPPTXバイト列に変換する

    Args:
        figs: plotly.graph_objects.Figure のリスト
        titles: 各スライドのタイトル（省略時はfig番号）

    Returns:
        PPTXファイルのバイト列
    """
    import tempfile

    if titles is None:
        titles = [f"Plot {i + 1}" for i in range(len(figs))]

    plot_images: list[dict[str, Any]] = []

    with tempfile.TemporaryDirectory() as tmpdir:
        for i, fig in enumerate(figs):
            img_path = Path(tmpdir) / f"plot_{i}.png"
            try:
                fig.write_image(str(img_path), width=1920, height=1080, scale=2)
                plot_images.append({"path": img_path, "title": titles[i] if i < len(titles) else f"Plot {i + 1}"})
            except Exception:
                # kaleido未インストール等の場合スキップ
                continue

        if not plot_images:
            raise ValueError("プロット画像の生成に失敗しました。kaleidoがインストールされているか確認してください。")

        return export_plots_to_pptx_bytes(plot_images)


# ================================================================
# ヘルパー関数
# ================================================================


def _add_slide_title(slide: Any, title: str) -> None:
    """スライドにタイトルテキストボックスを追加する"""
    from pptx.util import Inches

    left = Inches(MARGIN_INCH)
    top = Inches(MARGIN_INCH * 0.3)
    width = Inches(USABLE_WIDTH_INCH)
    height = Inches(TITLE_HEIGHT_INCH * 0.8)

    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.name = "メイリオ"


def _place_images_in_grid(
    slide: Any,
    images: list[Path],
    cols: int,
    rows: int,
) -> None:
    """スライド上に画像をグリッド配置する"""
    cell_width = USABLE_WIDTH_INCH / cols
    cell_height = USABLE_HEIGHT_INCH / rows

    for idx, img_path in enumerate(images):
        col = idx % cols
        row = idx // cols

        left = Inches(MARGIN_INCH + col * cell_width + CELL_PADDING_INCH)
        top = Inches(MARGIN_INCH + TITLE_HEIGHT_INCH + row * cell_height + CELL_PADDING_INCH)
        max_w = Inches(cell_width - 2 * CELL_PADDING_INCH)
        max_h = Inches(cell_height - 2 * CELL_PADDING_INCH)

        _add_image_fit(slide, img_path, left, top, max_w, max_h)


def _add_image_fit(
    slide: Any,
    img_path: Path,
    left: Emu,
    top: Emu,
    max_width: Emu,
    max_height: Emu,
) -> None:
    """画像をアスペクト比を保ったまま指定領域に収める"""
    from PIL import Image

    try:
        with Image.open(img_path) as img:
            img_w, img_h = img.size
    except Exception:
        # PIL不可の場合はそのまま配置
        slide.shapes.add_picture(str(img_path), left, top, max_width, max_height)
        return

    # アスペクト比を保ってフィット
    ratio_w = max_width / (img_w * 914400 / 96)  # px → EMU (96dpi想定)
    ratio_h = max_height / (img_h * 914400 / 96)
    ratio = min(ratio_w, ratio_h, 1.0)  # 拡大しない

    final_w = int(img_w * 914400 / 96 * ratio)
    final_h = int(img_h * 914400 / 96 * ratio)

    # 中央寄せ
    offset_x = (max_width - final_w) // 2
    offset_y = (max_height - final_h) // 2

    slide.shapes.add_picture(
        str(img_path),
        left + offset_x,
        top + offset_y,
        final_w,
        final_h,
    )
