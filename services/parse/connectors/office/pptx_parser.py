"""PPTXメタデータパーサー: .pptxファイルからメタデータを抽出

python-pptxライブラリを使用して以下の情報を抽出する:
- スライド数
- テキストコンテンツ（各スライドのタイトル・本文）
- 画像数
- スライドサイズ
- コアプロパティ（作成者・タイトル等）

[READMEへ戻る](../../../../README.md)
"""

from __future__ import annotations

import logging
import sys
from typing import TYPE_CHECKING, Any

from services.parse.base import AbstractFileParser

if TYPE_CHECKING:
    from services.graph.project_graph import ProjectGraph

logger = logging.getLogger(__name__)


class PptxMetadataParser(AbstractFileParser):
    """PPTXファイルのメタデータを抽出してノードプロパティに付与する

    priority=35: 入出力関係解析(30-33)の後、ディレクトリ関係(50)の前に実行
    """

    priority = 35

    def apply(self, graph: ProjectGraph) -> ProjectGraph:
        try:
            from pptx import Presentation
        except ImportError:
            logger.debug("python-pptx が未インストールのため PPTX パーサーをスキップ")
            return graph

        for node in graph.nodes:
            if node.format != "pptx":
                continue

            file_path = graph.project_root / node.properties.get("path", "")
            if not file_path.exists():
                continue

            try:
                metadata = _extract_pptx_metadata(file_path, Presentation)
                node.properties.update(metadata)
            except Exception as e:
                print(
                    f"警告: PPTX メタデータ抽出に失敗 ({file_path.name}): {e}",
                    file=sys.stderr,
                )

        return graph


def _extract_pptx_metadata(file_path: Any, Presentation: Any) -> dict[str, Any]:
    """PPTXファイルからメタデータを抽出する

    Args:
        file_path: PPTXファイルのパス
        Presentation: python-pptx の Presentation クラス

    Returns:
        抽出されたメタデータの辞書
    """
    prs = Presentation(str(file_path))
    metadata: dict[str, Any] = {}

    # スライド数
    slide_count = len(prs.slides)
    metadata["pptx.slide_count"] = slide_count

    # スライドサイズ（EMU → cm に変換）
    if prs.slide_width and prs.slide_height:
        metadata["pptx.slide_width_cm"] = round(prs.slide_width / 914400 * 2.54, 2)
        metadata["pptx.slide_height_cm"] = round(prs.slide_height / 914400 * 2.54, 2)

    # 各スライドのタイトルとテキスト、画像数
    slide_titles: list[str] = []
    total_image_count = 0
    total_text_length = 0

    for slide in prs.slides:
        # タイトル取得
        title = ""
        if slide.shapes.title is not None:
            title = slide.shapes.title.text.strip()
        slide_titles.append(title)

        # テキストとイメージのカウント
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    total_text_length += len(text)

            if shape.shape_type is not None:
                # shape_type 13 = PICTURE
                try:
                    if shape.shape_type == 13:
                        total_image_count += 1
                except Exception:
                    pass

    metadata["pptx.image_count"] = total_image_count
    metadata["pptx.total_text_length"] = total_text_length

    # スライドタイトル一覧（空でないもの）
    non_empty_titles = [t for t in slide_titles if t]
    if non_empty_titles:
        metadata["pptx.slide_titles"] = non_empty_titles

    # コアプロパティ
    core = prs.core_properties
    if core.author:
        metadata["pptx.author"] = core.author
    if core.title:
        metadata["pptx.title"] = core.title
    if core.subject:
        metadata["pptx.subject"] = core.subject
    if core.created:
        metadata["pptx.created"] = str(core.created)
    if core.modified:
        metadata["pptx.modified"] = str(core.modified)

    return metadata
